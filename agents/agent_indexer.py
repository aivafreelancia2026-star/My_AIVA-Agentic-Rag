"""
agents/agent_indexer.py

Indexes each agent's data sources (files + URLs) into its own FAISS vector store.
Each agent gets its own subfolder:
  agents/indexes/<agent_id>/faiss_index/
  agents/indexes/<agent_id>/chunks.json      ← raw text chunks for BM25 fallback
  agents/indexes/<agent_id>/_status.json     ← {"status": "indexing"|"indexed"|"failed", "error": "..."}
"""

import os
import json
import time
import logging
from pathlib import Path
from typing import List, Dict, Optional, Any

logger = logging.getLogger(__name__)

AGENTS_DIR   = Path(__file__).parent
INDEXES_DIR  = AGENTS_DIR / "indexes"


# ── Index-dir helpers ─────────────────────────────────────────────────────────

def _agent_index_dir(agent_id: str) -> Path:
    d = INDEXES_DIR / agent_id
    d.mkdir(parents=True, exist_ok=True)
    return d


def _write_status(agent_id: str, status: str, error: str = "") -> None:
    p = _agent_index_dir(agent_id) / "_status.json"
    with open(p, "w", encoding="utf-8") as f:
        json.dump({"status": status, "error": error, "ts": time.time()}, f)


def read_index_status(agent_id: str) -> Dict:
    """Return {"status": "...", "error": ""}. Status: pending|indexing|indexed|failed."""
    status_file = INDEXES_DIR / agent_id / "_status.json"
    faiss_dir   = INDEXES_DIR / agent_id / "faiss_index"

    if faiss_dir.exists():
        return {"status": "indexed", "error": ""}
    if status_file.exists():
        try:
            return json.loads(status_file.read_text(encoding="utf-8"))
        except Exception:
            pass
    return {"status": "pending", "error": ""}


# ── URL scraping ──────────────────────────────────────────────────────────────

def _scrape_url(url: str) -> str:
    try:
        import requests
        from html.parser import HTMLParser

        class _Stripper(HTMLParser):
            def __init__(self):
                super().__init__()
                self._parts = []
                self._skip = False
            def handle_starttag(self, tag, attrs):
                if tag in ("script", "style", "nav", "footer", "header"):
                    self._skip = True
            def handle_endtag(self, tag):
                if tag in ("script", "style", "nav", "footer", "header"):
                    self._skip = False
            def handle_data(self, data):
                if not self._skip:
                    s = data.strip()
                    if s:
                        self._parts.append(s)

        resp = requests.get(url, timeout=20,
                            headers={"User-Agent": "Mozilla/5.0 (AIVA-RAG/1.0)"})
        resp.raise_for_status()
        p = _Stripper()
        p.feed(resp.text)
        return "\n".join(p._parts)
    except Exception as e:
        logger.warning(f"[indexer] URL scrape failed for {url}: {e}")
        return ""


# ── File reading — supports all notebook file types ───────────────────────────

def _read_file(path: str) -> str:
    """
    Extract plain text from a file.
    Handles: PDF, DOCX/DOC, PPTX/PPT, XLSX, XLS, TXT, MD, CSV, JSON
    """
    p = Path(path)
    if not p.exists():
        logger.warning(f"[indexer] File not found: {path}")
        return ""

    suffix = p.suffix.lower()
    try:
        # ── PDF ──────────────────────────────────────────────────────────
        if suffix == ".pdf":
            try:
                import fitz  # PyMuPDF
                doc = fitz.open(str(p))
                pages = []
                for page in doc:
                    text = page.get_text("text")
                    if text.strip():
                        pages.append(text)
                doc.close()
                return "\n\n".join(pages)
            except ImportError:
                logger.error("[indexer] PyMuPDF (fitz) not installed — cannot read PDF")
                return ""

        # ── Word ─────────────────────────────────────────────────────────
        elif suffix in (".docx", ".doc"):
            try:
                from docx import Document as DocxDocument
                doc = DocxDocument(str(p))
                parts = []
                for para in doc.paragraphs:
                    if para.text.strip():
                        parts.append(para.text)
                # Also extract tables
                for table in doc.tables:
                    for row in table.rows:
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells:
                            parts.append(" | ".join(cells))
                return "\n\n".join(parts)
            except ImportError:
                logger.error("[indexer] python-docx not installed — cannot read DOCX")
                return ""

        # ── PowerPoint ───────────────────────────────────────────────────
        elif suffix in (".pptx", ".ppt"):
            try:
                from pptx import Presentation
                prs = Presentation(str(p))
                slides = []
                for i, slide in enumerate(prs.slides, 1):
                    texts = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            texts.append(shape.text.strip())
                    if texts:
                        slides.append(f"[Slide {i}]\n" + "\n".join(texts))
                return "\n\n".join(slides)
            except ImportError:
                logger.error("[indexer] python-pptx not installed — cannot read PPTX")
                return ""

        # ── Excel ─────────────────────────────────────────────────────────
        elif suffix == ".xlsx":
            try:
                import openpyxl
                wb = openpyxl.load_workbook(str(p), data_only=True, read_only=True)
                sheets = []
                for ws in wb.worksheets:
                    rows = []
                    for row in ws.iter_rows(values_only=True):
                        cells = [str(c) for c in row if c is not None and str(c).strip()]
                        if cells:
                            rows.append(" | ".join(cells))
                    if rows:
                        sheets.append(f"[Sheet: {ws.title}]\n" + "\n".join(rows))
                wb.close()
                return "\n\n".join(sheets)
            except ImportError:
                logger.error("[indexer] openpyxl not installed — cannot read XLSX")
                return ""

        elif suffix == ".xls":
            try:
                import xlrd
                wb = xlrd.open_workbook(str(p))
                sheets = []
                for ws in wb.sheets():
                    rows = []
                    for ri in range(ws.nrows):
                        cells = [str(ws.cell_value(ri, ci)).strip()
                                 for ci in range(ws.ncols)
                                 if str(ws.cell_value(ri, ci)).strip()]
                        if cells:
                            rows.append(" | ".join(cells))
                    if rows:
                        sheets.append(f"[Sheet: {ws.name}]\n" + "\n".join(rows))
                return "\n\n".join(sheets)
            except ImportError:
                logger.error("[indexer] xlrd not installed — cannot read XLS")
                return ""

        # ── Markdown ──────────────────────────────────────────────────────
        elif suffix in (".md", ".markdown"):
            raw = p.read_text(encoding="utf-8", errors="ignore")
            try:
                import markdown
                from html.parser import HTMLParser

                class _S(HTMLParser):
                    def __init__(self):
                        super().__init__()
                        self._d = []
                    def handle_data(self, d):
                        if d.strip():
                            self._d.append(d.strip())

                html = markdown.markdown(raw)
                s = _S()
                s.feed(html)
                return "\n".join(s._d)
            except ImportError:
                return raw  # Fall back to raw markdown text

        # ── Plain text / CSV / JSON ───────────────────────────────────────
        elif suffix == ".txt":
            return p.read_text(encoding="utf-8", errors="ignore")

        elif suffix == ".csv":
            import csv
            rows = []
            with open(p, encoding="utf-8", errors="ignore") as f:
                for row in csv.reader(f):
                    rows.append(", ".join(row))
            return "\n".join(rows)

        elif suffix == ".json":
            data = json.loads(p.read_text(encoding="utf-8"))
            return json.dumps(data, indent=2)

        else:
            # Last resort — try reading as UTF-8 text
            return p.read_text(encoding="utf-8", errors="ignore")

    except Exception as e:
        logger.warning(f"[indexer] File read error {path}: {e}")
        return ""


# ── Chunker ───────────────────────────────────────────────────────────────────

def _chunk_text(text: str, chunk_size: int = 800, overlap: int = 150) -> List[str]:
    if not text.strip():
        return []
    chunks = []
    start = 0
    while start < len(text):
        end   = start + chunk_size
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        start += chunk_size - overlap
    return chunks


# ── Core indexing ─────────────────────────────────────────────────────────────

def index_agent(agent: Dict, embedding_model) -> bool:
    """
    Build (or rebuild) the FAISS index for a single agent.
    Writes _status.json throughout so the UI can track progress.
    Returns True on success.
    """
    agent_id    = agent["id"]
    sources     = agent.get("sources", [])
    index_dir   = _agent_index_dir(agent_id)
    faiss_path  = index_dir / "faiss_index"
    chunks_path = index_dir / "chunks.json"

    _write_status(agent_id, "indexing")
    logger.info(f"[indexer] Indexing agent '{agent_id}' ({len(sources)} sources)...")

    all_texts: List[str] = []
    all_meta:  List[Dict] = []

    for source in sources:
        stype = source.get("type", "")
        sval  = source.get("value", "")

        if stype == "url":
            logger.info(f"[indexer]   Scraping URL: {sval}")
            raw   = _scrape_url(sval)
            label = sval
        elif stype == "file":
            logger.info(f"[indexer]   Reading file: {sval}")
            raw   = _read_file(sval)
            label = Path(sval).name
        else:
            logger.warning(f"[indexer]   Unknown source type: {stype}")
            continue

        if not raw.strip():
            logger.warning(f"[indexer]   ⚠️  No text extracted from: {sval}")
            continue

        chunks = _chunk_text(raw)
        logger.info(f"[indexer]   → {len(chunks)} chunks from '{label}'")
        for c in chunks:
            all_texts.append(c)
            all_meta.append({"source": label, "type": stype, "agent_id": agent_id})

    if not all_texts:
        msg = f"No text could be extracted from any source for agent '{agent_id}'"
        logger.warning(f"[indexer] ❌ {msg}")
        _write_status(agent_id, "failed", msg)
        return False

    try:
        from langchain_community.vectorstores import FAISS
        from langchain_core.documents import Document

        logger.info(f"[indexer] Embedding {len(all_texts)} chunks for '{agent_id}'...")
        docs = [
            Document(page_content=t, metadata=m)
            for t, m in zip(all_texts, all_meta)
        ]
        vectorstore = FAISS.from_documents(docs, embedding_model)
        vectorstore.save_local(str(faiss_path))

        # Save raw chunks for BM25 / keyword fallback
        with open(chunks_path, "w", encoding="utf-8") as f:
            json.dump(
                [{"text": t, "meta": m} for t, m in zip(all_texts, all_meta)],
                f, ensure_ascii=False, indent=2
            )

        _write_status(agent_id, "indexed")
        logger.info(f"[indexer] ✅ Agent '{agent_id}' indexed: {len(all_texts)} chunks.")
        return True

    except Exception as e:
        msg = str(e)
        logger.error(f"[indexer] ❌ FAISS indexing failed for '{agent_id}': {msg}")
        _write_status(agent_id, "failed", msg)
        return False


# ── Load / search ─────────────────────────────────────────────────────────────

def load_agent_index(agent_id: str, embedding_model):
    """Load the FAISS vectorstore for an agent. Returns None if not indexed yet."""
    try:
        from langchain_community.vectorstores import FAISS
        faiss_path = _agent_index_dir(agent_id) / "faiss_index"
        if not faiss_path.exists():
            return None
        return FAISS.load_local(
            str(faiss_path),
            embedding_model,
            allow_dangerous_deserialization=True,
        )
    except Exception as e:
        logger.warning(f"[indexer] Could not load index for '{agent_id}': {e}")
        return None


def agent_is_indexed(agent_id: str) -> bool:
    return (INDEXES_DIR / agent_id / "faiss_index").exists()


def search_agent_index(
    agent_id: str,
    query: str,
    embedding_model,
    k: int = 6,
) -> List[Dict]:
    """Search a single agent's FAISS index. Returns list of {text, source, score}."""
    vs = load_agent_index(agent_id, embedding_model)
    if vs is None:
        return []
    try:
        results = vs.similarity_search_with_score(query, k=k)
        return [
            {
                "text":     doc.page_content,
                "source":   doc.metadata.get("source", agent_id),
                "agent_id": agent_id,
                "score":    float(score),
            }
            for doc, score in results
        ]
    except Exception as e:
        logger.warning(f"[indexer] Search error for agent '{agent_id}': {e}")
        return []
