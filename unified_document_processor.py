# unified_document_processor.py - Comprehensive Multi-Format Document Processing System

# Consolidated processor supporting PDF and PowerPoint documents with advanced OCR, image analysis, and table extraction



import os

import re

import time

import json

import mimetypes

import hashlib

import tempfile

import subprocess

import io

import base64

from pathlib import Path

from typing import List, Dict, Any, Optional, Tuple, Union, Set

from datetime import datetime

from dataclasses import dataclass

from collections import defaultdict



# Core document types

from langchain_core.documents import Document


from langchain_text_splitters import RecursiveCharacterTextSplitter



# Import configuration

from config import *



# ---------- OCR AND IMAGE PROCESSING IMPORTS ----------



# OCR imports

try:

    import easyocr

    HAS_EASYOCR = True

except ImportError:

    HAS_EASYOCR = False



# PDF processing imports
try:
    import fitz  # PyMuPDF
    HAS_PYMUPDF = True
except ImportError:
    HAS_PYMUPDF = False
    print("Warning: PyMuPDF not available - PDF processing disabled")

# Table extraction imports
try:
    import tabula
    HAS_TABULA = True
except ImportError:
    HAS_TABULA = False
    print("Warning: tabula-py not available - table extraction limited")

try:
    import camelot
    HAS_CAMELOT = True
except ImportError:
    HAS_CAMELOT = False
    print("Warning: camelot not available - table extraction limited")

# Image processing imports for OCR preprocessing

try:

    from PIL import Image, ImageEnhance, ImageFilter, ImageOps

    HAS_PIL = True

except ImportError:

    HAS_PIL = False



# NLP imports for name extraction

try:

    import spacy

    HAS_SPACY = True

    # Load English model for name extraction

    try:

        nlp = spacy.load("en_core_web_sm")

    except OSError:

        print("Warning: spaCy English model not found. Install with: python -m spacy download en_core_web_sm")

        nlp = None

        HAS_SPACY = False

except ImportError:

    HAS_SPACY = False

    nlp = None



# ---------- ADDITIONAL PROCESSING IMPORTS ----------



# Import gracefully with fallbacks

try:

    import magic

    HAS_MAGIC = True

except ImportError:

    HAS_MAGIC = False

    print("Warning: python-magic not available - file type detection limited")



# Word document processing
try:
    import docx
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False
    print("Warning: python-docx not available - Word processing disabled")



try:

    from pptx import Presentation

    HAS_PPTX = True

except ImportError:

    HAS_PPTX = False

    print("Warning: python-pptx not available - PowerPoint processing disabled")



# Excel processing
try:
    import openpyxl
    HAS_EXCEL = True
except ImportError:
    HAS_EXCEL = False
    print("Warning: openpyxl not available - Excel processing disabled")

# CSV/text processing (always available via stdlib)
HAS_TEXT = True

# XML processing removed - only PDF and PPT supported

HAS_XML = False



# Audio processing removed - only PDF and PPT supported

HAS_AUDIO = False



# Archive processing removed - only PDF and PPT supported

HAS_ARCHIVE = False



# Global OCR reader (initialized once for efficiency)

_ocr_reader = None



def get_ocr_reader():

    """Get or initialize the OCR reader"""

    global _ocr_reader

    if _ocr_reader is None and HAS_EASYOCR:

        try:

            _ocr_reader = easyocr.Reader(['en'])

        except Exception as e:

            print(f"Warning: Could not initialize OCR reader: {e}")

            return None

    return _ocr_reader



# ---------- CORE DATA STRUCTURES ----------



@dataclass

class ProcessedDocument:

    """Represents a processed document with metadata"""

    content: str

    metadata: Dict[str, Any]

    doc_type: str

    file_path: str

    processing_time: float

    chunks: List[Document] = None



# ---------- OCR AND TEXT EXTRACTION FUNCTIONS ----------



def extract_ui_elements_from_image(image_path: str) -> Dict[str, Any]:
    """Extract UI elements (popups, dialogs, buttons, screens) from image screenshots"""
    if not HAS_EASYOCR:
        return {'text': '', 'ui_elements': {}}
    
    reader = get_ocr_reader()
    if not reader:
        return {'text': '', 'ui_elements': {}}
    
    try:
        processed_path = _preprocess_image_for_ocr(image_path)
        
        # Get detailed OCR results with bounding boxes
        results = reader.readtext(processed_path, 
                                detail=1, 
                                paragraph=False,
                                width_ths=0.5,
                                height_ths=0.5,
                                text_threshold=0.3,
                                low_text=0.2,
                                link_threshold=0.2)
        
        # Analyze UI elements based on position and text patterns
        ui_elements = {
            'dialog_titles': [],
            'popup_names': [],
            'screen_names': [],
            'buttons': [],
            'tabs': [],
            'labels': [],
            'values': [],
            'all_text': []
        }
        
        all_text_items = []
        
        for (bbox, text, confidence) in results:
            text = text.strip()
            if not text or confidence < 0.3:
                continue
            
            text = _clean_ocr_text(text)
            if not text:
                continue
            
            # Get position information
            top_left = bbox[0]
            y_position = top_left[1]
            x_position = top_left[0]
            
            text_item = {
                'text': text,
                'confidence': confidence,
                'y_pos': y_position,
                'x_pos': x_position,
                'bbox': bbox
            }
            all_text_items.append(text_item)
            ui_elements['all_text'].append(text)
            
            # Detect UI element types based on patterns
            text_lower = text.lower()
            
            # Dialog/Popup titles (usually at top, longer text, high confidence)
            if (confidence > 0.6 and len(text) > 10 and 
                ('snapshot' in text_lower or 'performance' in text_lower or 
                 'report' in text_lower or 'dashboard' in text_lower or
                 '-' in text or ':' in text)):
                ui_elements['dialog_titles'].append(text)
                ui_elements['popup_names'].append(text)
            
            # Screen/Page names
            if ('page' in text_lower or 'screen' in text_lower or 
                'trend' in text_lower or 'view' in text_lower or
                'appears' in text_lower):
                ui_elements['screen_names'].append(text)
            
            # Buttons (short text, common button words)
            button_words = ['ok', 'cancel', 'next', 'previous', 'apply', 'close', 
                          'save', 'submit', 'drill down', 'export', 'print']
            if (len(text) <= 20 and any(btn in text_lower for btn in button_words)):
                ui_elements['buttons'].append(text)
            
            # Tabs (short text, often single words or abbreviations)
            tab_words = ['qa', 'coaching', 'parameters', 'frequency', 'score', 
                        'low scoring', 'overview', 'details', 'summary']
            if len(text) <= 15 and any(tab in text_lower for tab in tab_words):
                ui_elements['tabs'].append(text)
            
            # Labels (text ending with :)
            if text.endswith(':') or (len(text) < 30 and confidence > 0.5):
                ui_elements['labels'].append(text)
            
            # Values (numbers, percentages, dates)
            if re.match(r'^\d+[\d\.\,\%\-\/]*$', text):
                ui_elements['values'].append(text)
        
        # Sort text by position (top to bottom, left to right) for structured output
        all_text_items.sort(key=lambda x: (x['y_pos'], x['x_pos']))
        
        # Create structured text output
        structured_text = []
        
        if ui_elements['dialog_titles']:
            structured_text.append(f"DIALOG/POPUP: {' | '.join(set(ui_elements['dialog_titles']))}")
        
        if ui_elements['screen_names']:
            structured_text.append(f"SCREEN/PAGE: {' | '.join(set(ui_elements['screen_names']))}")
        
        if ui_elements['tabs']:
            structured_text.append(f"TABS: {', '.join(set(ui_elements['tabs']))}")
        
        if ui_elements['buttons']:
            structured_text.append(f"BUTTONS: {', '.join(set(ui_elements['buttons']))}")
        
        # Add all text in reading order
        structured_text.append("ALL TEXT: " + " | ".join([item['text'] for item in all_text_items]))
        
        full_text = "\n".join(structured_text)
        
        return {
            'text': full_text,
            'simple_text': " ".join(ui_elements['all_text']),
            'ui_elements': ui_elements,
            'text_items': all_text_items
        }
        
    except Exception as e:
        print(f"OCR error for {image_path}: {e}")
        return {'text': '', 'ui_elements': {}}
    
    finally:
        if processed_path != image_path and os.path.exists(processed_path):
            try:
                os.remove(processed_path)
            except:
                pass


def extract_text_from_image(image_path: str) -> str:

    """Extract text from image using OCR with enhanced settings and preprocessing"""

    # Use the new enhanced UI element extraction
    result = extract_ui_elements_from_image(image_path)
    return result.get('text', '') if result.get('text') else result.get('simple_text', '')



def _preprocess_image_for_ocr(image_path: str) -> str:

    """Preprocess image for better OCR results"""

    if not HAS_PIL:

        return image_path

    

    try:

        # Create a temporary processed image path

        processed_path = image_path.replace('.png', '_ocr_processed.png').replace('.jpg', '_ocr_processed.jpg')

        

        with Image.open(image_path) as img:

            # Convert to RGB if needed

            if img.mode != 'RGB':

                img = img.convert('RGB')

            

            # Resize if too large (OCR works better on reasonable sizes)

            max_size = 2000

            if img.width > max_size or img.height > max_size:

                img.thumbnail((max_size, max_size), Image.Resampling.LANCZOS)

            

            # Convert to grayscale for better text detection

            gray_img = ImageOps.grayscale(img)

            

            # Enhance contrast

            enhancer = ImageEnhance.Contrast(gray_img)

            enhanced_img = enhancer.enhance(1.3)

            

            # Enhance sharpness

            enhancer = ImageEnhance.Sharpness(enhanced_img)

            sharpened_img = enhancer.enhance(1.2)

            

            # Apply noise reduction

            denoised_img = sharpened_img.filter(ImageFilter.MedianFilter(size=3))

            

            # Save processed image

            denoised_img.save(processed_path, 'PNG', optimize=True)

            

            return processed_path

            

    except Exception as e:

       # print(f"Image preprocessing failed: {e}")

        return image_path



def _is_navigation_button(ocr_text: str) -> bool:

    """Check if OCR text indicates this is a navigation button (previous/next)"""

    if not ocr_text:

        return False

    

    # Convert to lowercase for case-insensitive matching

    text_lower = ocr_text.lower().strip()

    

    # Navigation button patterns - exact word matches only

    nav_patterns = [

        'previous',

        'next',

        'prev',

        'back',

        'forward',

        'continue',

        'proceed',

        'skip',

        'close',

        'exit',

        'done',

        'finish',

        'submit',

        'ok',

        'cancel',

        'yes',

        'no'

    ]

    

    # Split text into words for exact matching

    words = text_lower.split()

    

    # Check if any word exactly matches navigation patterns

    for word in words:

        # Clean the word (remove punctuation)

        clean_word = word.strip('.,!?;:()[]{}"\'')

        if clean_word in nav_patterns:

            # Only detect as navigation button if it's:

            # 1. A single word, OR

            # 2. Two words where the second is a common button word (like "go back", "click next")

            if len(words) == 1:

                return True

            elif len(words) == 2:

                # Check if the second word is also a navigation word or common button word

                second_word = words[1].strip('.,!?;:()[]{}"\'')

                button_words = ['go', 'click', 'press', 'tap', 'button', 'link']

                if second_word in button_words or second_word in nav_patterns:

                    return True

    

    return False



def _clean_ocr_text(text: str) -> str:

    """Clean common OCR errors and improve text quality"""

    if not text:

        return ""

    

    # Common OCR error patterns and fixes

    fixes = {

        # Ubuntu/command line fixes

        'ubuntudubuntu': 'ubuntu@ubuntu',

        'ubuntuqubuntu': 'ubuntu@ubuntu',

        'ubuntudubuntu:': 'ubuntu@ubuntu:',

        'ubuntuqubuntu:': 'ubuntu@ubuntu:',

        'lakshithadlakshitha': 'lakshitha@lakshitha',

        'lakshithadlakshitha:': 'lakshitha@lakshitha:',

        

        # Directory/path fixes

        'Desktops': 'Desktop',

        'desktop': 'Desktop',

        'nedias': 'media',

        'nynewshare': 'mynewshare',

        'nyapps': 'myapp',

        'nyapp': 'myapp',

        

        # Command line fixes

        'sudo apt install': 'sudo apt install',

        'sudo apt-get install': 'sudo apt-get install',

        'Reading package lists': 'Reading package lists',

        'Building dependency tree': 'Building dependency tree',

        'sstdio.hz': 'stdio.h',

        'hello worldln': 'hello world',

        'hello.€': 'hello.c',

        'hello.â‚¬': 'hello.c',

        'requirments.txt': 'requirements.txt',

        'requirnents': 'requirements',

        

        # Common character substitutions

        '€': 'c',  # Euro symbol often misread as 'c'

        'â‚¬': 'c',  # Another euro symbol variant

        'z': 'z',  # Keep as is

        

        # Programming fixes

        'printf("hello worldln");': 'printf("hello world");',

        'fron flask': 'from flask',

        'inport Flask': 'import Flask',

        'apachez': 'apache2',

        'systenctl': 'systemctl',

        'Py thon3': 'python3',

        'Imyenv': 'myenv',

        'Imynenv': 'myenv',

        

        # File extension fixes

        '.€': '.c',

        '.â‚¬': '.c',

        '.Py': '.py',

        '.py': '.py',

    }

    

    cleaned_text = text

    

    # Apply fixes

    for error, fix in fixes.items():

        cleaned_text = cleaned_text.replace(error, fix)

    

    # Fix common patterns with regex

    # Fix repeated characters (common OCR error)

    cleaned_text = re.sub(r'(\w)\1{2,}', r'\1', cleaned_text)

    

    # Fix common command line patterns

    cleaned_text = re.sub(r'bash: cd: ([^:]+): No such file or directory', r'bash: cd: \1: No such file or directory', cleaned_text)

    

    # Remove excessive punctuation but keep important ones

    cleaned_text = re.sub(r'[^\w\s\n\.\,\:\;\!\?\(\)\[\]\-\+\=\%\$\@\#\/\\\"]', ' ', cleaned_text)

    

    # Clean up multiple spaces and newlines

    cleaned_text = re.sub(r'\s+', ' ', cleaned_text)

    cleaned_text = re.sub(r'\n\s*\n', '\n', cleaned_text)

    

    return cleaned_text.strip()



def extract_names_from_content(content: str) -> List[str]:

    """Extract person names from document content using spaCy"""

    if not HAS_SPACY or not nlp or not content.strip():

        return []

    

    try:

        doc = nlp(content)

        names = []

        

        # Extract named entities that are persons

        for ent in doc.ents:

            if ent.label_ == "PERSON":

                # Clean up the name

                name = ent.text.strip()

                # Filter out single words and very short names

                if len(name.split()) >= 2 and len(name) > 3:

                    names.append(name)

        

        # Remove duplicates while preserving order

        unique_names = []

        seen = set()

        for name in names:

            if name.lower() not in seen:

                unique_names.append(name)

                seen.add(name.lower())

        

        return unique_names[:10]  # Limit to top 10 names

        

    except Exception as e:

        print(f"Name extraction error: {e}")

        return []



# ---------- FILE TYPE DETECTION ----------



class FileTypeDetector:

    """Advanced file type detection and validation"""

    

    SUPPORTED_FORMATS = {

        # Text documents

        'pdf': ['.pdf'],
        'docx': ['.docx', '.doc'],

        # Plain text / markup
        'text': ['.txt', '.md', '.csv'],

        # Spreadsheets
        'xlsx': ['.xlsx', '.xls'],

        # Presentations

        'pptx': ['.pptx', '.ppt']

    }

    

    @classmethod

    def get_file_type(cls, file_path: str) -> Tuple[str, str]:

        """

        Detect file type and return (category, specific_type)

        Returns: ('pdf', 'pdf') or ('video', 'mp4') etc.

        """

        file_path = Path(file_path)

        extension = file_path.suffix.lower()

        

        # Try magic library first if available

        if HAS_MAGIC:

            try:

                mime_type = magic.from_file(str(file_path), mime=True)

                return cls._mime_to_category(mime_type, extension)

            except Exception:

                pass

        

        # Fallback to extension-based detection

        for category, extensions in cls.SUPPORTED_FORMATS.items():

            if extension in extensions:

                return category, extension[1:]  # Remove dot

        

        return 'unknown', extension[1:] if extension else 'unknown'

    

    @classmethod

    def _mime_to_category(cls, mime_type: str, extension: str) -> Tuple[str, str]:

        """Convert MIME type to our category system"""

        mime_mappings = {

            'application/pdf': ('pdf', 'pdf'),

            'application/vnd.openxmlformats-officedocument.presentationml.presentation': ('pptx', 'pptx'),

            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': ('docx', 'docx'),

            'application/msword': ('docx', 'doc'),

            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': ('xlsx', 'xlsx'),

            'application/vnd.ms-excel': ('xlsx', 'xls'),

            'text/plain': ('text', 'txt'),

            'text/csv': ('text', 'csv'),

            'text/markdown': ('text', 'md'),

        }

        

        if mime_type in mime_mappings:

            return mime_mappings[mime_type]

        

        # Fallback to extension

        specific_type = extension[1:] if extension.startswith('.') else extension

        return 'unknown', specific_type

    

    @classmethod

    def is_supported(cls, file_path: str) -> bool:

        """Check if file type is supported"""

        category, _ = cls.get_file_type(file_path)

        return category != 'unknown'

    

    @classmethod

    def get_supported_extensions(cls) -> List[str]:

        """Get list of all supported file extensions"""

        extensions = []

        for ext_list in cls.SUPPORTED_FORMATS.values():

            extensions.extend(ext_list)

        return sorted(extensions)



# ---------- GRAPH PROCESSOR ----------



class GraphProcessor:

    """Advanced graph and chart analysis processor"""

    

    def __init__(self):

        self.chart_types = {

            'bar_chart': ['bar', 'column', 'histogram'],

            'line_chart': ['line', 'trend', 'time series'],

            'pie_chart': ['pie', 'donut', 'circle'],

            'scatter_plot': ['scatter', 'plot', 'correlation'],

            'flow_chart': ['flow', 'process', 'workflow', 'diagram'],

            'organizational': ['org', 'hierarchy', 'structure'],

            'network': ['network', 'graph', 'node', 'connection']

        }

        

        self.chart_keywords = [

            'chart', 'graph', 'plot', 'figure', 'diagram', 'visualization',

            'data', 'trend', 'analysis', 'statistics', 'metrics', 'performance'

        ]

    

    def analyze_chart_content(self, image_info: Dict, ocr_text: str) -> Dict:

        """Enhanced chart content analysis"""

        analysis = {

            'chart_type': 'unknown',

            'data_points': [],

            'trends': [],

            'insights': [],

            'numerical_data': [],

            'labels': [],

            'title': '',

            'axes_info': {}

        }

        

        if not ocr_text:

            return analysis

        

        text_lower = ocr_text.lower()

        

        # Determine chart type

        for chart_type, keywords in self.chart_types.items():

            if any(keyword in text_lower for keyword in keywords):

                analysis['chart_type'] = chart_type

                break

        

        # Extract numerical data

        numbers = re.findall(r'\d+(?:\.\d+)?(?:[kmb%])?', ocr_text, re.IGNORECASE)

        analysis['numerical_data'] = numbers

        

        # Extract potential labels

        lines = ocr_text.split('\n')

        potential_labels = []

        for line in lines:

            line = line.strip()

            if line and not line.isdigit() and len(line) < 50:

                potential_labels.append(line)

        analysis['labels'] = potential_labels

        

        # Look for title

        title_indicators = ['title', 'chart', 'graph', 'figure']

        for line in lines[:3] + lines[-2:]:

            if any(indicator in line.lower() for indicator in title_indicators):

                analysis['title'] = line.strip()

                break

        

        # Extract trends and insights

        trend_keywords = ['increase', 'decrease', 'growth', 'decline', 'stable', 'peak', 'valley']

        for keyword in trend_keywords:

            if keyword in text_lower:

                analysis['trends'].append(keyword)

        

        # Detect axes information

        axes_keywords = ['x-axis', 'y-axis', 'horizontal', 'vertical', 'time', 'value', 'category']

        for keyword in axes_keywords:

            if keyword in text_lower:

                analysis['axes_info'][keyword] = True

        

        return analysis

    

    def generate_chart_description(self, analysis: Dict, context: str = "") -> str:

        """Generate comprehensive chart description for embedding"""

        description_parts = [

            f"CHART ANALYSIS - Type: {analysis['chart_type'].upper()}"

        ]

        

        if analysis['title']:

            description_parts.append(f"Title: {analysis['title']}")

        

        if analysis['numerical_data']:

            description_parts.append(f"Data Values: {', '.join(analysis['numerical_data'][:10])}")

        

        if analysis['labels']:

            description_parts.append(f"Labels/Categories: {', '.join(analysis['labels'][:8])}")

        

        if analysis['trends']:

            description_parts.append(f"Identified Trends: {', '.join(analysis['trends'])}")

        

        if analysis['axes_info']:

            description_parts.append(f"Chart Elements: {', '.join(analysis['axes_info'].keys())}")

        

        if context:

            description_parts.append(f"Context: {context}")

        

        description_parts.append(f"Chart Type Keywords: {', '.join(self.chart_types.get(analysis['chart_type'], []))}")

        description_parts.append("This visualization contains quantitative data and can be referenced for data analysis questions.")

        

        return '\n'.join(description_parts)



# ---------- ENHANCED IMAGE PROCESSOR ----------



class EnhancedImageProcessor:

    """Advanced image processor with OCR and analysis"""

    

    def __init__(self):

        self.ocr_reader = None

        self.processed_images: Set[str] = set()

        self.graph_processor = GraphProcessor()

        

        if HAS_EASYOCR:

            try:

                self.ocr_reader = easyocr.Reader(['en'], gpu=False)

                print("OCR reader initialized successfully")

            except Exception as e:

                print(f"OCR initialization failed: {e}")

    

    def _calculate_image_hash(self, image_data: bytes) -> str:

        """Calculate hash for deduplication"""

        return hashlib.md5(image_data).hexdigest()

    

    def _preprocess_image(self, image_path: str) -> str:

        """Preprocess image for better OCR results"""

        if not HAS_PIL:

            return image_path

        

        try:

            with Image.open(image_path) as img:

                if img.mode != 'RGB':

                    img = img.convert('RGB')

                

                if img.width > config.IMAGE_MAX_SIZE or img.height > config.IMAGE_MAX_SIZE:

                    img.thumbnail((config.IMAGE_MAX_SIZE, config.IMAGE_MAX_SIZE), Image.Resampling.LANCZOS)

                

                enhancer = ImageEnhance.Contrast(img)

                img = enhancer.enhance(1.2)

                

                enhancer = ImageEnhance.Sharpness(img)

                img = enhancer.enhance(1.1)

                

                preprocessed_path = image_path.replace('.', '_processed.')

                img.save(preprocessed_path, 'JPEG', quality=95)

                

                return preprocessed_path

        

        except Exception as e:

            #  print(f"Image preprocessing failed: {e}")

            return image_path

    

    def extract_images_with_context(self, pdf_path: str, output_dir: str) -> List[Dict]:

        """Extract images with enhanced context and analysis"""

        if not HAS_PYMUPDF:

            print("PyMuPDF not available - skipping image extraction")

            return []

        

        images_info = []

        os.makedirs(output_dir, exist_ok=True)

        self.processed_images.clear()

        

        try:

            print(f"Opening PDF: {pdf_path}")

            doc = fitz.open(pdf_path)

            total_pages = len(doc)

            print(f"Total pages to process: {total_pages}")

            

            total_images_found = 0

            total_images_processed = 0

            

            for page_num in range(total_pages):

                try:

                    print(f"Processing page {page_num + 1}/{total_pages}...")

                    page = doc.load_page(page_num)

                    

                    page_text = page.get_text()

                    image_list = page.get_images(full=True)

                    page_image_count = len(image_list)

                    total_images_found += page_image_count

                    

                    if page_image_count > 0:

                        print(f"  Found {page_image_count} images on page {page_num + 1}")

                    

                    for img_index, img_info in enumerate(image_list):

                        try:

                            xref = img_info[0]

                            img_dict = doc.extract_image(xref)

                            image_data = img_dict["image"]

                            

                            # Deduplication

                            img_hash = self._calculate_image_hash(image_data)

                            if img_hash in self.processed_images:

                                print(f"    Skipping duplicate image {img_index}")

                                continue

                            

                            # Size filtering

                            pix = fitz.Pixmap(doc, xref)

                            if (pix.width < config.IMAGE_MIN_SIZE or pix.height < config.IMAGE_MIN_SIZE or 

                                (pix.width > 2500 and pix.height > 2500)):

                                print(f"    Skipping image {img_index} due to size ({pix.width}x{pix.height})")

                                pix = None

                                continue

                            

                            print(f"    Processing image {img_index}: {pix.width}x{pix.height}")

                            

                            # Enhanced context extraction

                            context_text = ""

                            structured_context = {}

                            

                            try:

                                img_rects = page.get_image_rects(xref)

                                if img_rects:

                                    img_rect = img_rects[0]

                                    

                                    zones = {

                                        'above': fitz.Rect(img_rect.x0, max(0, img_rect.y0 - 100), 

                                                         img_rect.x1, img_rect.y0),

                                        'below': fitz.Rect(img_rect.x0, img_rect.y1, 

                                                         img_rect.x1, min(page.rect.height, img_rect.y1 + 100)),

                                        'left': fitz.Rect(max(0, img_rect.x0 - 150), img_rect.y0, 

                                                        img_rect.x0, img_rect.y1),

                                        'right': fitz.Rect(img_rect.x1, img_rect.y0, 

                                                         min(page.rect.width, img_rect.x1 + 150), img_rect.y1)

                                    }

                                    

                                    for zone_name, zone_rect in zones.items():

                                        zone_text = page.get_text(clip=zone_rect).strip()

                                        if zone_text:

                                            structured_context[zone_name] = zone_text

                                    

                                    context_text = " ".join(structured_context.values())

                            

                            except Exception as e:

                                print(f"      Context extraction failed: {e}")

                            

                            # Save image

                            pdf_name = Path(pdf_path).stem

                            img_filename = f"{pdf_name}_page{page_num+1}_img{img_index}.{img_dict['ext']}"

                            img_path = os.path.join(output_dir, img_filename)

                            

                            try:

                                with open(img_path, "wb") as f:

                                    f.write(image_data)

                                print(f"    Saved: {img_filename}")

                            except Exception as e:

                                print(f"    Failed to save image: {e}")

                                pix = None

                                continue

                            

                            # Preprocess for better OCR

                            preprocessed_path = self._preprocess_image(img_path)

                            

                            # Enhanced OCR

                            ocr_text = ""

                            if self.ocr_reader:

                                try:

                                    ocr_text = self.perform_enhanced_ocr(preprocessed_path)

                                    if ocr_text:

                                        print(f"    OCR extracted {len(ocr_text)} characters")

                                except Exception as e:

                                    print(f"    OCR failed: {e}")

                                

                                # Clean up preprocessed file if different

                                if preprocessed_path != img_path:

                                    try:

                                        os.remove(preprocessed_path)

                                    except:

                                        pass

                            

                            # Skip logos

                            if self.is_company_logo(ocr_text, pix.width, pix.height):

                                print(f"    Detected as logo, removing...")

                                try:

                                    os.remove(img_path)

                                except:

                                    pass

                                pix = None

                                continue

                            

                            # Advanced image analysis

                            image_type = self.classify_image_content(ocr_text, pix.width, pix.height, context_text)

                            

                            # Enhanced chart analysis

                            chart_analysis = {}

                            if image_type in ['chart', 'graph', 'diagram']:

                                chart_analysis = self.graph_processor.analyze_chart_content(

                                    {'width': pix.width, 'height': pix.height}, ocr_text

                                )

                            

                            # Create comprehensive image info

                            image_info = {

                                "path": img_path,

                                "filename": img_filename,

                                "pdf_source": os.path.basename(pdf_path),

                                "page": page_num,

                                "display_page": page_num + 1,

                                "image_index": img_index,  # Use actual image index from enumerate

                                "width": pix.width,

                                "height": pix.height,

                                "aspect_ratio": pix.width / pix.height if pix.height > 0 else 1.0,

                                "format": img_dict.get('ext', 'unknown'),

                                "colorspace": img_dict.get('colorspace', 'unknown'),

                                "xref": xref,

                                "hash": img_hash,

                                "ocr_text": ocr_text,

                                "surrounding_context": context_text.strip(),

                                "structured_context": structured_context,

                                "page_text_snippet": self._extract_page_snippet(page_text, 400),

                                "image_type": image_type,

                                "chart_analysis": chart_analysis,

                                "extraction_timestamp": datetime.now().isoformat()

                            }

                            

                            images_info.append(image_info)

                            self.processed_images.add(img_hash)

                            total_images_processed += 1

                            

                            pix = None

                            

                        except Exception as e:

                            print(f"    Error processing image {img_index}: {e}")

                            continue

                

                except Exception as e:

                    print(f"  Error processing page {page_num + 1}: {e}")

                    continue

            

            doc.close()

            

            print(f"Image extraction complete:")

            print(f"  Total images found: {total_images_found}")

            print(f"  Images processed: {total_images_processed}")

            print(f"  Images saved: {len(images_info)}")

            

            # Save metadata

            if images_info:

                metadata_path = os.path.join(output_dir, "images_metadata.json")

                try:

                    with open(metadata_path, 'w', encoding='utf-8') as f:

                        json.dump(images_info, f, indent=2, ensure_ascii=False)

                    print(f"  Metadata saved: {metadata_path}")

                except Exception as e:

                    print(f"  Failed to save metadata: {e}")

            

            return images_info

            

        except Exception as e:

            print(f"Critical error in image extraction: {e}")

            return []

    

    def perform_enhanced_ocr(self, image_path: str) -> str:

        """Enhanced OCR with better accuracy"""

        if not self.ocr_reader:

            return ""

        

        try:

            results = self.ocr_reader.readtext(image_path, 

                                             paragraph=False,

                                             width_ths=0.7,

                                             height_ths=0.7)

            

            if not results:

                return ""

            

            processed_results = []

            

            for (bbox, text, confidence) in results:

                if confidence > config.OCR_CONFIDENCE_THRESHOLD and text.strip():

                    cleaned_text = re.sub(r'\s+', ' ', text.strip())

                    

                    if len(cleaned_text) > 1 or confidence > 0.8:

                        processed_results.append((bbox, cleaned_text, confidence))

            

            # Sort by position (top to bottom, left to right)

            processed_results.sort(key=lambda x: (x[0][0][1], x[0][0][0]))

            

            # Group into lines based on y-coordinate proximity

            lines = []

            current_line = []

            current_y = None

            

            for bbox, text, confidence in processed_results:

                y_pos = bbox[0][1]

                

                if current_y is None or abs(y_pos - current_y) < 20:

                    current_line.append((bbox[0][0], text))

                    current_y = y_pos if current_y is None else current_y

                else:

                    if current_line:

                        current_line.sort(key=lambda x: x[0])

                        line_text = " ".join([item[1] for item in current_line])

                        lines.append(line_text)

                    

                    current_line = [(bbox[0][0], text)]

                    current_y = y_pos

            

            if current_line:

                current_line.sort(key=lambda x: x[0])

                line_text = " ".join([item[1] for item in current_line])

                lines.append(line_text)

            

            final_text = "\n".join(lines)

            final_text = re.sub(r'\n\s*\n', '\n', final_text)

            final_text = re.sub(r'[^\w\s\n\.\,\:\;\!\?\(\)\[\]\-\+\=\%\$\@\#]', ' ', final_text)

            

            return final_text.strip()

            

        except Exception as e:

            print(f"Enhanced OCR error for {image_path}: {e}")

            return ""

    

    def classify_image_content(self, ocr_text: str, width: int, height: int, context: str = "") -> str:

        """Enhanced image classification"""

        combined_text = f"{ocr_text} {context}".lower()

        aspect_ratio = width / height if height > 0 else 1.0

        

        classification_keywords = {

            'chart': ['chart', 'graph', 'plot', 'data', 'axis', 'trend', 'statistics', 'metrics'],

            'table': ['table', 'row', 'column', 'header', 'name', 'title', 'list', 'signature', 'department'],

            'diagram': ['diagram', 'flow', 'process', 'workflow', 'step', 'flowchart', 'procedure'],

            'equation': ['equation', 'formula', '=', 'calculate', 'sum', '+', '-', '×', '÷', 'math'],

            'organizational': ['org', 'hierarchy', 'structure', 'manager', 'director', 'team'],

            'network': ['network', 'connection', 'node', 'link', 'relationship']

        }

        

        scores = {}

        for category, keywords in classification_keywords.items():

            score = sum(1 for keyword in keywords if keyword in combined_text)

            

            # Adjust score based on context

            if category == 'chart' and any(chart_word in combined_text for chart_word in ['bar', 'line', 'pie']):

                score += 2

            elif category == 'table' and aspect_ratio > 1.5:  # Wide images often tables

                score += 1

            elif category == 'diagram' and 1.0 <= aspect_ratio <= 2.0:  # Square-ish diagrams

                score += 1

            

            scores[category] = score

        

        # Determine best classification

        if scores:

            best_category = max(scores, key=scores.get)

            if scores[best_category] > 0:

                return best_category

        

        # Fallback classification based on size and aspect ratio

        if aspect_ratio > 2.5:

            return 'table'

        elif 0.5 <= aspect_ratio <= 1.5:

            return 'chart'

        else:

            return 'image'

    

    def is_company_logo(self, ocr_text: str, width: int, height: int) -> bool:

        """Enhanced logo detection"""

        if not ocr_text:

            return False

        

        ocr_lower = ocr_text.lower().strip()

        

        # Size constraints

        is_small = width < 300 and height < 300

        is_very_small = width < 120 or height < 120

        

        # Text constraints

        has_minimal_text = len(ocr_text.strip()) < 40

        word_count = len(ocr_text.split())

        

        # Company indicators

        has_company_name = any(variant in ocr_lower for variant in config.COMPANY_VARIANTS)

        

        # Logo indicators

        logo_indicators = [

            'logo', 'copyright', '©', 'confidential', 'proprietary',

            'trademark', '®', '™', 'all rights reserved', 'inc', 'llc', 'ltd'

        ]

        has_logo_words = any(indicator in ocr_lower for indicator in logo_indicators)

        

        # Decision logic with stricter criteria

        if is_very_small and (has_company_name or has_logo_words):

            return True

        

        if has_company_name and is_small and word_count < 5:

            return True

        

        if has_logo_words and has_minimal_text and word_count < 6:

            return True

        

        # Additional check for very short text that's likely logo text

        if word_count <= 2 and is_small and len(ocr_text.strip()) < 20:

            return True

        

        return False

    

    def _extract_page_snippet(self, page_text: str, max_length: int = 400) -> str:

        """Extract meaningful snippet from page text"""

        if not page_text:

            return ""

        

        # Clean the text

        cleaned = re.sub(r'\s+', ' ', page_text.strip())

        

        # Try to find the most informative section

        paragraphs = cleaned.split('\n\n')

        if paragraphs:

            # Find paragraph with good balance of content

            best_paragraph = max(paragraphs, 

                               key=lambda p: len(p) if 50 < len(p) < max_length * 2 else 0)

            

            if len(best_paragraph) <= max_length:

                return best_paragraph

        

        # Fallback: middle section

        if len(cleaned) > max_length:

            mid_point = len(cleaned) // 2

            start = max(0, mid_point - max_length // 2)

            end = min(len(cleaned), start + max_length)

        return cleaned



# ---------- ADVANCED TABLE PROCESSOR ----------



class AdvancedTableProcessor:

    """Advanced table processing with multiple extraction methods"""

    

    def __init__(self):

        self.table_patterns = {

            'pipe_table': r'\|[^|]*\|[^|]*\|',

            'box_drawing': r'[┌─└┘├┤┬┴┼│]',

            'multi_column': r'^\s*\w+\s{2,}\w+\s{2,}\w+',

            'header_separator': r'^[-=_]{3,}',

            'tabular_data': r'^\s*[^\s]+\s{2,}[^\s]+\s{2,}[^\s]+',

            'csv_like': r'[^,\n]+,[^,\n]+,[^,\n]+',

            'whitespace_delimited': r'^\s*\S+\s+\S+\s+\S+.*'

        }

        

        self.table_headers = [

            'name', 'title', 'role', 'position', 'department', 'email', 'phone', 

            'signature', 'date', 'status', 'id', 'number', 'code', 'value', 

            'amount', 'quantity', 'price', 'description', 'category', 'type'

        ]

        

        self.strong_table_indicators = [

            r'table\s+\d+', r'figure\s+\d+', r'appendix\s+[a-z]',

            r'list\s+of', r'directory', r'roster', r'staff', r'employee'

        ]

    

    def extract_tables_with_tabula(self, pdf_path: str, page_num: int) -> List:

        """Extract tables using tabula-py for better accuracy"""

        if not HAS_TABULA:

            return []

        

        try:

            methods = ['stream', 'lattice']

            all_tables = []

            

            for method in methods:

                try:

                    tables = tabula.read_pdf(

                        pdf_path,

                        pages=page_num + 1,

                        multiple_tables=True,

                        pandas_options={'header': 0},

                        stream=(method == 'stream'),

                        lattice=(method == 'lattice'),

                        silent=True

                    )

                    

                    if tables:

                        all_tables.extend(tables)

                        print(f"    Tabula ({method}) found {len(tables)} tables")

                

                except Exception:

                    continue

            

            # Filter and clean tables

            cleaned_tables = []

            for table in all_tables:

                if hasattr(table, 'dropna') and not table.empty:

                    table = table.dropna(how='all')

                    table = table.loc[:, ~table.columns.str.contains('^Unnamed')]

                    

                    if table.shape[0] >= config.TABLE_MIN_ROWS and table.shape[1] >= 2:

                        cleaned_tables.append(table)

            

            return cleaned_tables

            

        except Exception as e:

            print(f"    Tabula extraction failed: {e}")

            return []

    

    def extract_tables_with_camelot(self, pdf_path: str, page_num: int) -> List:

        """Extract tables using camelot for complex tables"""

        if not HAS_CAMELOT:

            return []

        

        try:

            tables = []

            methods = ['stream', 'lattice']

            

            for method in methods:

                try:

                    camelot_tables = camelot.read_pdf(

                        pdf_path,

                        pages=str(page_num + 1),

                        flavor=method,

                        suppress_stdout=True

                    )

                    

                    if camelot_tables:

                        for table in camelot_tables:

                            if hasattr(table, 'df') and table.df is not None and not table.df.empty:

                                if hasattr(table, 'accuracy') and table.accuracy > 50:

                                    tables.append(table.df)

                                elif not hasattr(table, 'accuracy'):

                                    tables.append(table.df)

                        

                        print(f"    Camelot ({method}) found {len(camelot_tables)} tables")

                

                except Exception:

                    continue

            

            return tables

            

        except Exception as e:

            print(f"    Camelot extraction failed: {e}")

            return []

    

    def detect_tables_with_context(self, text: str, page_num: int = 0, pdf_path: str = "") -> List[Dict]:

        """Enhanced table detection with multiple methods"""

        tables = []

        lines = text.split('\n')

        

        # Method 1: Enhanced pipe table detection

        tables.extend(self._detect_pipe_tables(lines))

        

        # Method 2: Whitespace-delimited tables

        tables.extend(self._detect_whitespace_tables(lines))

        

        # Method 3: Pattern-based detection

        tables.extend(self._detect_pattern_tables(lines))

        

        # Method 4: External library extraction (if available)

        if pdf_path and (HAS_TABULA or HAS_CAMELOT):

            external_tables = []

            

            if HAS_TABULA:

                tabula_tables = self.extract_tables_with_tabula(pdf_path, page_num)

                external_tables.extend(tabula_tables)

            

            if HAS_CAMELOT:

                camelot_tables = self.extract_tables_with_camelot(pdf_path, page_num)

                external_tables.extend(camelot_tables)

            

            # Convert external tables to our format

            for i, df_table in enumerate(external_tables):

                table_info = {

                    'type': 'external_extraction',

                    'start_line': 0,

                    'end_line': len(lines),

                    'content': df_table.to_string(index=False),

                    'dataframe': df_table,

                    'pre_context': '',

                    'post_context': '',

                    'full_context': df_table.to_string(index=False),

                    'extraction_method': 'tabula/camelot',

                    'confidence': 0.9

                }

                tables.append(table_info)

        

        # Enhanced deduplication and ranking

        tables = self._deduplicate_and_rank_tables(tables)

        

        return tables

    

    def _detect_pipe_tables(self, lines: List[str]) -> List[Dict]:

        """Enhanced pipe table detection"""

        tables = []

        table_start = None

        table_lines = []

        pre_context = []

        

        for i, line in enumerate(lines):

            line_stripped = line.strip()

            

            if ('|' in line and line.count('|') >= 2 and 

                not line_stripped.startswith('|') or line_stripped.endswith('|')):

                

                if table_start is None:

                    table_start = i

                    pre_context = lines[max(0, i-3):i]

                table_lines.append(line)

            else:

                if table_start is not None and len(table_lines) >= config.TABLE_MIN_ROWS:

                    post_context = lines[i:min(len(lines), i+3)]

                    

                    if self._validate_table_quality(table_lines):

                        table_info = {

                            'type': 'pipe_table',

                            'start_line': table_start,

                            'end_line': i-1,

                            'content': '\n'.join(table_lines),

                            'pre_context': '\n'.join(pre_context),

                            'post_context': '\n'.join(post_context),

                            'full_context': '\n'.join(pre_context + table_lines + post_context),

                            'confidence': 0.8

                        }

                        tables.append(table_info)

                

                table_start = None

                table_lines = []

                pre_context = []

        

        return tables

    

    def _detect_whitespace_tables(self, lines: List[str]) -> List[Dict]:

        """Detect tables separated by whitespace"""

        tables = []

        potential_table_lines = []

        table_start = None

        

        for i, line in enumerate(lines):

            line_stripped = line.strip()

            

            if self._is_potential_table_row(line_stripped):

                if table_start is None:

                    table_start = i

                potential_table_lines.append((i, line))

            else:

                if len(potential_table_lines) >= config.TABLE_MIN_ROWS:

                    if self._validate_whitespace_table(potential_table_lines):

                        pre_context = lines[max(0, table_start-2):table_start]

                        post_context = lines[i:min(len(lines), i+2)]

                        

                        table_content = '\n'.join([line for _, line in potential_table_lines])

                        

                        table_info = {

                            'type': 'whitespace_table',

                            'start_line': table_start,

                            'end_line': potential_table_lines[-1][0],

                            'content': table_content,

                            'pre_context': '\n'.join(pre_context),

                            'post_context': '\n'.join(post_context),

                            'full_context': '\n'.join(pre_context + [table_content] + post_context),

                            'confidence': 0.7

                        }

                        tables.append(table_info)

                

                potential_table_lines = []

                table_start = None

        

        return tables

    

    def _detect_pattern_tables(self, lines: List[str]) -> List[Dict]:

        """Detect tables using various patterns"""

        tables = []

        

        for i in range(len(lines) - 2):

            current_line = lines[i].strip()

            next_line = lines[i+1].strip()

            

            if (current_line and next_line and

                self._looks_like_header(current_line) and

                self._looks_like_separator(next_line)):

                

                table_lines = [current_line, next_line]

                table_end = i + 1

                

                for j in range(i + 2, min(i + 20, len(lines))):

                    data_line = lines[j].strip()

                    if not data_line:

                        break

                    if self._looks_like_data_row(data_line, current_line):

                        table_lines.append(data_line)

                        table_end = j

                    else:

                        break

                

                if len(table_lines) >= 3:

                    pre_context = lines[max(0, i-2):i]

                    post_context = lines[table_end+1:min(len(lines), table_end+3)]

                    

                    table_info = {

                        'type': 'header_separator_table',

                        'start_line': i,

                        'end_line': table_end,

                        'content': '\n'.join(table_lines),

                        'pre_context': '\n'.join(pre_context),

                        'post_context': '\n'.join(post_context),

                        'full_context': '\n'.join(pre_context + table_lines + post_context),

                        'confidence': 0.75

                    }

                    tables.append(table_info)

        

        return tables

    

    def _validate_table_quality(self, table_lines: List[str]) -> bool:

        """Validate if detected table is of good quality"""

        if len(table_lines) < config.TABLE_MIN_ROWS:

            return False

        

        pipe_counts = [line.count('|') for line in table_lines]

        if len(set(pipe_counts)) > 2:

            return False

        

        total_chars = sum(len(line) for line in table_lines)

        if total_chars < 50:

            return False

        

        return True

    

    def _is_potential_table_row(self, line: str) -> bool:

        """Check if line could be part of a whitespace-delimited table"""

        if not line or len(line) < 10:

            return False

        

        parts = re.split(r'\s{2,}', line)

        

        if len(parts) >= 2:

            non_empty_parts = [p for p in parts if p.strip()]

            return len(non_empty_parts) >= 2

        

        return False

    

    def _validate_whitespace_table(self, potential_lines: List[Tuple[int, str]]) -> bool:

        """Validate whitespace-separated table"""

        if len(potential_lines) < config.TABLE_MIN_ROWS:

            return False

        

        column_counts = []

        for _, line in potential_lines:

            parts = re.split(r'\s{2,}', line.strip())

            non_empty_parts = [p for p in parts if p.strip()]

            column_counts.append(len(non_empty_parts))

        

        if len(set(column_counts)) > 3:

            return False

        

        combined_text = ' '.join([line for _, line in potential_lines]).lower()

        table_score = sum(1 for header in self.table_headers if header in combined_text)

        

        return table_score > 0

    

    def _looks_like_header(self, line: str) -> bool:

        """Check if line looks like a table header"""

        header_indicators = ['name', 'title', 'id', 'date', 'status', 'type', 'description']

        line_lower = line.lower()

        

        header_count = sum(1 for indicator in header_indicators if indicator in line_lower)

        

        parts = re.split(r'\s{2,}|\t|\|', line)

        clean_parts = [p.strip() for p in parts if p.strip()]

        

        return header_count > 0 and len(clean_parts) >= 2

    

    def _looks_like_separator(self, line: str) -> bool:

        """Check if line looks like a separator"""

        if not line:

            return False

        

        separator_chars = set('-=_')

        line_chars = set(line.replace(' ', ''))

        

        if line_chars <= separator_chars and len(line) > 5:

            return True

        

        separator_count = len(re.findall(r'[-=_]{2,}', line))

        return separator_count >= 2

    

    def _looks_like_data_row(self, line: str, header_line: str) -> bool:

        """Check if line looks like a data row for given header"""

        header_parts = len(re.split(r'\s{2,}|\t|\|', header_line))

        data_parts = len(re.split(r'\s{2,}|\t|\|', line))

        

        return abs(header_parts - data_parts) <= 1

    

    def _deduplicate_and_rank_tables(self, tables: List[Dict]) -> List[Dict]:

        """Remove duplicates and rank tables by quality"""

        if not tables:

            return []

        

        unique_tables = []

        seen_content = set()

        

        for table in tables:

            content_hash = hashlib.md5(table['content'].encode()).hexdigest()

            if content_hash not in seen_content:

                seen_content.add(content_hash)

                unique_tables.append(table)

        

        def table_quality_score(table):

            confidence = table.get('confidence', 0.5)

            content_length = len(table['content'])

            line_count = table['content'].count('\n') + 1

            

            return confidence * 100 + content_length * 0.1 + line_count * 5

        

        unique_tables.sort(key=table_quality_score, reverse=True)

        

        return unique_tables[:5]

    

    def extract_structured_data(self, table_content: str) -> Dict:

        """Enhanced structured data extraction"""

        lines = [line.strip() for line in table_content.split('\n') if line.strip()]

        

        if not lines:

            return {'headers': [], 'data_rows': [], 'row_count': 0, 'column_count': 0}

        

        headers = []

        data_rows = []

        

        if '|' in table_content:

            for i, line in enumerate(lines):

                if '|' in line:

                    cells = [cell.strip() for cell in line.split('|')]

                    cells = [cell for cell in cells if cell]

                    

                    if cells:

                        if (i == 0 or not headers) and self._looks_like_header(' '.join(cells)):

                            headers = cells

                        elif not (set(line) <= set('-=_|')):

                            data_rows.append(cells)

        

        else:

            for i, line in enumerate(lines):

                if set(line.replace(' ', '')) <= set('-=_'):

                    continue

                

                parts = re.split(r'\s{2,}|\t', line)

                clean_parts = [part.strip() for part in parts if part.strip()]

                

                if clean_parts:

                    if (i < 2 or not headers) and self._looks_like_header(' '.join(clean_parts)):

                        headers = clean_parts

                    else:

                        data_rows.append(clean_parts)

        

        # Post-processing: ensure consistent column count

        if headers and data_rows:

            target_columns = len(headers)

            filtered_rows = []

            

            for row in data_rows:

                if abs(len(row) - target_columns) <= 1:

                    if len(row) < target_columns:

                        row.extend([''] * (target_columns - len(row)))

                    elif len(row) > target_columns:

                        row = row[:target_columns]

                    filtered_rows.append(row)

            

            data_rows = filtered_rows

        

        return {

            'headers': headers,

            'data_rows': data_rows,

            'row_count': len(data_rows),

            'column_count': len(headers) if headers else (len(data_rows[0]) if data_rows else 0),

            'extraction_confidence': len(headers) > 0 and len(data_rows) > 0

        }

    

    def enhance_table_for_embedding(self, table_info: Dict) -> str:

        """Create enhanced table representation for better search"""

        structured_data = self.extract_structured_data(table_info['content'])

        

        enhanced_parts = [

            "[ENHANCED TABLE DATA - OPTIMIZED FOR SEARCH]",

            f"Table Type: {table_info['type']}",

            f"Confidence: {table_info.get('confidence', 'N/A')}",

            f"Extraction Method: {table_info.get('extraction_method', 'pattern_based')}"

        ]

        

        if table_info.get('pre_context'):

            enhanced_parts.append(f"Context Before Table: {table_info['pre_context']}")

        

        if structured_data['headers']:

            enhanced_parts.append(f"Table Headers: {' | '.join(structured_data['headers'])}")

            enhanced_parts.append(f"Column Names for Search: {', '.join(structured_data['headers'])}")

        

        enhanced_parts.append("Raw Table Content:")

        enhanced_parts.append(table_info['content'])

        

        if structured_data['headers'] and structured_data['data_rows']:

            enhanced_parts.append("Searchable Data Entries:")

            

            for row_idx, row in enumerate(structured_data['data_rows'][:15]):

                if len(row) >= len(structured_data['headers']):

                    row_entries = []

                    for header, value in zip(structured_data['headers'], row):

                        if header and value and str(value).strip():

                            row_entries.append(f"{header}: {value}")

                    

                    if row_entries:

                        enhanced_parts.append(f"Row {row_idx + 1}: {' | '.join(row_entries)}")

        

        enhanced_parts.append(f"Table Statistics: {structured_data['row_count']} rows, {structured_data['column_count']} columns")

        

        if table_info.get('post_context'):

            enhanced_parts.append(f"Context After Table: {table_info['post_context']}")

        

        search_tags = []

        if 'name' in ' '.join(structured_data.get('headers', [])).lower():

            search_tags.append("contains names and personal information")

        if any(title in ' '.join(structured_data.get('headers', [])).lower() 

               for title in ['title', 'position', 'role']):

            search_tags.append("contains job titles and roles")

        if 'department' in ' '.join(structured_data.get('headers', [])).lower():

            search_tags.append("contains organizational structure")

        

        if search_tags:

            enhanced_parts.append(f"Table Contains: {', '.join(search_tags)}")

        

        enhanced_parts.append("[END ENHANCED TABLE DATA]")

        

        return '\n'.join(enhanced_parts)



# ---------- DOCUMENT PROCESSORS ----------



# WordProcessor removed - only PDF and PPT supported



# ---------- PDF PROCESSOR ----------

class PDFProcessor:
    """Process PDF documents with text extraction, image analysis, and table detection"""
    
    @staticmethod
    def process(file_path: str, output_dir: str = None) -> ProcessedDocument:
        """Process PDF document and extract content including images and tables"""
        start_time = time.time()
        
        if not HAS_PYMUPDF:
            raise Exception("PyMuPDF not available - cannot process PDF files")
        
        try:
            print(f"📄 Processing PDF: {Path(file_path).name}")
            
            # Open PDF document
            doc = fitz.open(file_path)
            total_pages = len(doc)
            
            # Initialize content and metadata
            full_content = []
            images_info = []
            tables_info = []
            page_texts = []
            
            # Process each page
            for page_num in range(total_pages):
                try:
                    page = doc.load_page(page_num)
                    
                    # Extract text from page
                    page_text = page.get_text()
                    if page_text.strip():
                        page_texts.append(f"Page {page_num + 1}:\n{page_text}\n")
                        full_content.append(f"Page {page_num + 1}:\n{page_text}\n")
                    
                    # Extract images if output directory is provided
                    if output_dir:
                        try:
                            image_processor = EnhancedImageProcessor()
                            page_images = image_processor.extract_images_with_context(file_path, output_dir)
                            images_info.extend(page_images)
                        except Exception as e:
                            print(f"Warning: Image extraction failed for page {page_num + 1}: {e}")
                    
                    # Extract tables if available
                    try:
                        table_processor = AdvancedTableProcessor()
                        page_tables = table_processor.detect_tables_with_context(page_text, page_num, file_path)
                        tables_info.extend(page_tables)
                        
                        # Add table content to the main content
                        for table_info in page_tables:
                            if 'content' in table_info and table_info['content']:
                                table_content = f"\n\n[TABLE on Page {page_num + 1}]:\n{table_info['content']}\n"
                                full_content.append(table_content)
                                
                    except Exception as e:
                        print(f"Warning: Table extraction failed for page {page_num + 1}: {e}")
                        
                except Exception as e:
                    print(f"Warning: Failed to process page {page_num + 1}: {e}")
                    continue
            
            doc.close()
            
            # Combine all content
            combined_content = "\n".join(full_content)
            
            # Create metadata
            processing_time = time.time() - start_time
            metadata = {
                'total_pages': total_pages,
                'word_count': len(combined_content.split()),
                'char_count': len(combined_content),
                'images_found': len(images_info),
                'tables_found': len(tables_info),
                'file_type': 'pdf',
                'processing_time': processing_time,
                'images_info': images_info,
                'tables_info': tables_info
            }
            
            # Extract document name and author heuristically
            doc_name = get_document_name(file_path)
            metadata['title'] = doc_name
            metadata['author'] = 'Unknown'  # Could be enhanced with PDF metadata extraction
            
            print(f"✅ PDF processed successfully:")
            print(f"   Pages: {total_pages}")
            print(f"   Content length: {len(combined_content)} chars")
            print(f"   Images found: {len(images_info)}")
            print(f"   Tables found: {len(tables_info)}")
            
            return ProcessedDocument(
                content=combined_content,
                metadata=metadata,
                doc_type='pdf',
                file_path=file_path,
                processing_time=processing_time
            )
            
        except Exception as e:
            raise Exception(f"Failed to process PDF {file_path}: {e}")

class PowerPointProcessor:

    """Process PowerPoint presentations (.pptx, .ppt) with enhanced segment-wise extraction"""

    

    @staticmethod
    def _extract_shape_text_detailed(shape, shape_idx: int, slide_num: int) -> Dict[str, Any]:
        """Extract text from shapes with detailed metadata including position and type"""
        result = {
            'text': '',
            'shape_type': str(shape.shape_type) if hasattr(shape, 'shape_type') else 'unknown',
            'shape_name': shape.name if hasattr(shape, 'name') else f'shape_{shape_idx}',
            'position': {},
            'is_text_box': False,
            'is_callout': False,
            'is_annotation': False
        }
        
        # Get position information
        if hasattr(shape, 'left') and hasattr(shape, 'top'):
            result['position'] = {
                'left': shape.left,
                'top': shape.top,
                'width': shape.width if hasattr(shape, 'width') else 0,
                'height': shape.height if hasattr(shape, 'height') else 0
            }
        
        # Check if it's a text box or callout
        if hasattr(shape, 'text') and shape.text.strip():
            result['text'] = shape.text.strip()
            
            # Detect if it's likely a callout/annotation (small text box, specific keywords)
            text_lower = result['text'].lower()
            annotation_indicators = ['note:', 'important:', 'note-', 'remember:', 'tip:', 'see', 'click', 'appears', 'shown', 'displayed']
            result['is_annotation'] = any(indicator in text_lower for indicator in annotation_indicators)
            
            # Check if it's a text box shape type
            if hasattr(shape, 'shape_type'):
                # TEXT_BOX = 17, PLACEHOLDER = 14
                result['is_text_box'] = str(shape.shape_type) in ['TEXT_BOX (17)', 'PLACEHOLDER (14)']
        
        # Extract text from shape's text frame
        if hasattr(shape, 'text_frame'):
            try:
                paragraphs = []
                for paragraph in shape.text_frame.paragraphs:
                    para_text = paragraph.text.strip()
                    if para_text:
                        paragraphs.append(para_text)
                if paragraphs and not result['text']:
                    result['text'] = '\n'.join(paragraphs)
            except:
                pass
        
        return result

    @staticmethod

    def process(file_path: str, output_dir: str = None) -> ProcessedDocument:

        """Process PowerPoint presentation with enhanced segment-wise extraction"""

        if not HAS_PPTX:

            raise ImportError("python-pptx required for PowerPoint processing")

        

        start_time = time.time()

        file_path = Path(file_path)

        

        try:

            prs = Presentation(str(file_path))

            

            content_parts = []

            slide_count = 0

            images_info = []
            all_slide_relationships = []  # Track relationships between slides

            

            # Create images directory

            if output_dir:

                images_dir = Path(output_dir) / "images"

            else:

                images_dir = file_path.parent / "images" / file_path.stem

            images_dir.mkdir(parents=True, exist_ok=True)

            

            for i, slide in enumerate(prs.slides, 1):

                slide_content = [f"\n{'='*80}"]
                slide_content.append(f"SLIDE {i} of {len(prs.slides)}")
                slide_content.append(f"{'='*80}")
                
                slide_text_parts = []
                slide_annotations = []  # Separate annotations from main content
                slide_callouts = []  # Callouts and explanatory text
                slide_image_count = 0
                
                # Extract speaker notes first
                if hasattr(slide, 'notes_slide') and slide.notes_slide:
                    try:
                        notes_text_frame = slide.notes_slide.notes_text_frame
                        if notes_text_frame and hasattr(notes_text_frame, 'text'):
                            notes = notes_text_frame.text.strip()
                            if notes:
                                slide_content.append(f"\n[SPEAKER NOTES]:")
                                slide_content.append(notes)
                                slide_content.append("")
                    except:
                        pass

                

                # Extract text, tables, and images from shapes with enhanced metadata
                shapes_data = []  # Store all shape data for relationship analysis
                
                for shape_idx, shape in enumerate(slide.shapes):
                    
                    # Extract detailed shape information
                    shape_details = PowerPointProcessor._extract_shape_text_detailed(shape, shape_idx, i)
                    shapes_data.append(shape_details)

                    # Separate main content from annotations
                    if shape_details['text']:
                        if shape_details['is_annotation']:
                            slide_annotations.append(f"[ANNOTATION]: {shape_details['text']}")
                        elif 'appears' in shape_details['text'].lower() or 'page' in shape_details['text'].lower():
                            slide_callouts.append(f"[CALLOUT on Slide {i}]: {shape_details['text']}")
                        else:
                            slide_text_parts.append(shape_details['text'])
                            slide_content.append(shape_details['text'])

                    

                    # Extract tables
                    if hasattr(shape, "table"):
                        try:
                            table = shape.table
                            table_content = []
                            
                            # Extract table headers
                            if table.rows:
                                header_row = table.rows[0]
                                headers = []
                                for cell in header_row.cells:
                                    headers.append(cell.text.strip())
                                table_content.append(" | ".join(headers))
                                table_content.append("-" * len(" | ".join(headers)))
                                
                                # Extract table data
                                for row in table.rows[1:]:
                                    row_data = []
                                    for cell in row.cells:
                                        row_data.append(cell.text.strip())
                                    table_content.append(" | ".join(row_data))
                            
                            if table_content:
                                table_text = "\n".join(table_content)
                                slide_content.append(f"\n[TABLE on Slide {i}]:\n{table_text}\n")
                                print(f"    Extracted table from slide {i}")
                                
                        except Exception as e:
                            print(f"    Warning: Failed to extract table from slide {i}: {e}")
                    
                    # Extract images

                    if hasattr(shape, 'image') and shape.image:

                        try:

                            image = shape.image

                            slide_image_count += 1  # Increment per-slide counter

                            image_filename = f"{file_path.stem}_slide{i}_img{shape_idx}.{image.ext}"

                            image_path = images_dir / image_filename

                            

                            # Save image

                            with open(image_path, 'wb') as f:

                                f.write(image.blob)

                            

                            # Extract text and UI elements from image using enhanced OCR

                            ocr_text = ""
                            ui_elements_data = {}

                            if HAS_EASYOCR:

                                print(f"    Running enhanced OCR on {image_filename}...")

                                ocr_result = extract_ui_elements_from_image(str(image_path))
                                ocr_text = ocr_result.get('text', '')
                                ui_elements_data = ocr_result.get('ui_elements', {})

                                if ocr_text:

                                    print(f"    OCR extracted: {len(ocr_text)} characters")
                                    
                                    # Show detected UI elements
                                    if ui_elements_data.get('dialog_titles'):
                                        print(f"       └─ Detected Popup/Dialog: {ui_elements_data['dialog_titles'][0]}")
                                    if ui_elements_data.get('screen_names'):
                                        print(f"       └─ Detected Screen: {ui_elements_data['screen_names'][0]}")
                                    if ui_elements_data.get('tabs'):
                                        print(f"       └─ Detected Tabs: {', '.join(ui_elements_data['tabs'][:3])}")
                                    if ui_elements_data.get('buttons'):
                                        print(f"       └─ Detected Buttons: {', '.join(ui_elements_data['buttons'][:3])}")

                            

                            # Check if this is a navigation button and skip if so

                            if _is_navigation_button(ocr_text):

                                print(f"    Skipping navigation button: {image_filename} (OCR: '{ocr_text}')")

                                # Clean up the saved image file

                                try:

                                    os.remove(image_path)

                                except:

                                    pass

                                continue

                            

                            # Get slide context including annotations and callouts

                            slide_context = " ".join(slide_text_parts)
                            annotations_context = " | ".join(slide_annotations) if slide_annotations else ""
                            callouts_context = " | ".join(slide_callouts) if slide_callouts else ""

                            

                            image_info = {

                                'path': str(image_path),

                                'filename': image_filename,

                                'slide_number': i,

                                'image_number_in_slide': shape_idx,

                                'shape_index': shape_idx,

                                'surrounding_context': slide_context,
                                'annotations': annotations_context,
                                'callouts': callouts_context,

                                'content_type': image.content_type,

                                'size': len(image.blob),

                                'extraction_method': 'pptx_embedded',

                                'ocr_text': ocr_text,

                                'has_ocr': bool(ocr_text),
                                'shape_details': shape_details,
                                
                                # Enhanced UI element detection
                                'ui_elements': ui_elements_data,
                                'popup_names': ui_elements_data.get('popup_names', []),
                                'dialog_titles': ui_elements_data.get('dialog_titles', []),
                                'screen_names': ui_elements_data.get('screen_names', []),
                                'detected_tabs': ui_elements_data.get('tabs', []),
                                'detected_buttons': ui_elements_data.get('buttons', []),
                                'has_popup': bool(ui_elements_data.get('popup_names', [])),
                                'has_dialog': bool(ui_elements_data.get('dialog_titles', [])),
                                'is_screenshot': bool(ui_elements_data.get('buttons', []) or ui_elements_data.get('tabs', []))

                            }

                            images_info.append(image_info)

                            

                            # Add image reference to content with full context including UI elements
                            context_parts = [f"[IMAGE {shape_idx} from SLIDE {i}]"]
                            context_parts.append(f"File: {image_filename}")
                            
                            # Add UI element detection results
                            if ui_elements_data.get('popup_names'):
                                context_parts.append(f"POPUP/DIALOG: {' | '.join(ui_elements_data['popup_names'])}")
                            if ui_elements_data.get('screen_names'):
                                context_parts.append(f"SCREEN: {' | '.join(ui_elements_data['screen_names'])}")
                            if ui_elements_data.get('tabs'):
                                context_parts.append(f"TABS: {', '.join(ui_elements_data['tabs'])}")
                            if ui_elements_data.get('buttons'):
                                context_parts.append(f"BUTTONS: {', '.join(ui_elements_data['buttons'])}")
                            
                            if slide_context:
                                context_parts.append(f"Main Context: {slide_context[:200]}")
                            if annotations_context:
                                context_parts.append(f"Annotations: {annotations_context}")
                            if callouts_context:
                                context_parts.append(f"Callouts: {callouts_context}")
                            if ocr_text:
                                context_parts.append(f"OCR Text: {ocr_text[:200]}")
                            
                            slide_content.append("\n" + " | ".join(context_parts))

                            

                            print(f"    Extracted image {shape_idx} from slide {i}: {image_filename}")

                            

                        except Exception as e:

                            print(f"    Warning: Could not extract image from slide {i}: {e}")

                    

                    # Handle tables in slides

                    if hasattr(shape, 'has_table') and shape.has_table:

                        table = shape.table

                        table_data = []

                        for row in table.rows:

                            row_data = [cell.text.strip() for cell in row.cells]

                            if any(row_data):

                                table_data.append(" | ".join(row_data))

                        

                        if table_data:

                            slide_content.append("TABLE:")

                            slide_content.extend(table_data)

                            slide_content.append("END TABLE")

                
                # Add annotations section if present
                if slide_annotations:
                    slide_content.append(f"\n[ANNOTATIONS & CALLOUTS on Slide {i}]:")
                    slide_content.extend(slide_annotations)
                
                # Add callouts section if present
                if slide_callouts:
                    slide_content.extend(slide_callouts)
                
                # Track slide relationships for navigation
                slide_relationship = {
                    'slide_number': i,
                    'has_navigation': any('previous' in text.lower() or 'next' in text.lower() 
                                         for text in slide_text_parts + slide_annotations + slide_callouts),
                    'references_other_slides': any('slide' in text.lower() or 'page' in text.lower() 
                                                   for text in slide_text_parts + slide_annotations + slide_callouts),
                    'text_objects_count': len(slide_text_parts),
                    'annotations_count': len(slide_annotations),
                    'callouts_count': len(slide_callouts),
                    'images_count': slide_image_count
                }
                all_slide_relationships.append(slide_relationship)

                

                if len(slide_content) > 1:  # More than just slide number

                    content_parts.extend(slide_content)

                    content_parts.append("")  # Empty line between slides

                    slide_count += 1

            

            content = "\n".join(content_parts)

            

            # Extract names from content

            content_names = extract_names_from_content(content)

            primary_author = content_names[0] if content_names else 'Unknown'

            

            # Extract metadata

            props = prs.core_properties

            # Calculate summary statistics
            total_annotations = sum(rel['annotations_count'] for rel in all_slide_relationships)
            total_callouts = sum(rel['callouts_count'] for rel in all_slide_relationships)
            slides_with_navigation = sum(1 for rel in all_slide_relationships if rel['has_navigation'])
            
            metadata = {

                'source': file_path.name,

                'file_type': 'pptx',

                'title': file_path.stem,  # Use actual file name

                'author': primary_author,  # Use content-based author

                'document_author': props.author or 'Unknown',  # Keep original document author

                'content_authors': content_names,  # All names found in content

                'subject': props.subject or '',

                'keywords': props.keywords or '',

                'created': props.created.isoformat() if props.created else '',

                'modified': props.modified.isoformat() if props.modified else '',

                'slide_count': slide_count,

                'total_slides': len(prs.slides),

                'image_count': len(images_info),

                'images': images_info,

                'has_images': len(images_info) > 0,

                'word_count': len(content.split()) if content else 0,
                
                # Enhanced segment-wise metadata
                'slide_relationships': all_slide_relationships,
                'total_annotations': total_annotations,
                'total_callouts': total_callouts,
                'slides_with_navigation': slides_with_navigation,
                'has_speaker_notes': any('SPEAKER NOTES' in content for _ in [1]),
                'extraction_enhanced': True,

                'processing_timestamp': datetime.now().isoformat()

            }

            

            processing_time = time.time() - start_time

            

            return ProcessedDocument(

                content=content,

                metadata=metadata,

                doc_type='pptx',

                file_path=str(file_path),

                processing_time=processing_time

            )

            

        except Exception as e:

            raise Exception(f"Failed to process PowerPoint presentation {file_path}: {e}")



# ExcelProcessor removed - only PDF and PPT supported



# XMLProcessor removed - only PDF and PPT supported



# AudioProcessor removed - only PDF and PPT supported



# TextProcessor removed - only PDF and PPT supported



# ---------- WORD PROCESSOR ----------

class WordProcessor:
    """Process Microsoft Word documents (.docx)"""

    @staticmethod
    def process(file_path: str, output_dir: str = None) -> ProcessedDocument:
        """Process Word document and extract text"""
        if not HAS_DOCX:
            raise ImportError("python-docx is not installed. Word processing is disabled.")

        start_time = time.time()
        file_path_obj = Path(file_path)
        
        try:
            doc = docx.Document(str(file_path_obj))
            
            # Extract text from paragraphs
            full_text = []
            for para in doc.paragraphs:
                if para.text.strip():
                    full_text.append(para.text.strip())
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_data = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_data.append(cell.text.strip())
                    if row_data:
                        full_text.append(" | ".join(row_data))
            
            content = "\n\n".join(full_text)
            
            # Metadata
            metadata = {
                'title': file_path_obj.stem,
                'author': 'Unknown',
                'word_count': len(content.split()),
                'page_count': 'N/A', # python-docx doesn't easily give page counts
                'creation_date': datetime.fromtimestamp(file_path_obj.stat().st_ctime).isoformat()
            }
            
            return ProcessedDocument(
                content=content,
                metadata=metadata,
                doc_type='docx',
                file_path=str(file_path_obj),
                processing_time=time.time() - start_time
            )
            
        except Exception as e:
            raise Exception(f"Word processing failed: {e}")


# ---------- TEXT / CSV PROCESSOR ----------

class TextProcessor:
    """Process plain text, markdown, and CSV files"""

    @staticmethod
    def process(file_path: str, output_dir: str = None) -> ProcessedDocument:
        start_time = time.time()
        file_path_obj = Path(file_path)
        ext = file_path_obj.suffix.lower()

        try:
            raw = file_path_obj.read_text(encoding='utf-8', errors='replace')

            if ext == '.csv':
                import csv as csv_mod
                import io
                reader = csv_mod.reader(io.StringIO(raw))
                rows = list(reader)
                lines = [" | ".join(r) for r in rows if any(c.strip() for c in r)]
                content = "\n".join(lines)
                doc_type = 'csv'
            else:
                content = raw
                doc_type = ext.lstrip('.')

            metadata = {
                'title': file_path_obj.stem,
                'author': 'Unknown',
                'word_count': len(content.split()),
                'char_count': len(content),
                'file_type': doc_type,
                'creation_date': datetime.fromtimestamp(file_path_obj.stat().st_ctime).isoformat()
            }

            return ProcessedDocument(
                content=content,
                metadata=metadata,
                doc_type=doc_type,
                file_path=str(file_path_obj),
                processing_time=time.time() - start_time
            )

        except Exception as e:
            raise Exception(f"Text processing failed for {file_path}: {e}")


# ---------- EXCEL PROCESSOR ----------

class ExcelProcessor:
    """Process Excel spreadsheets (.xlsx, .xls)"""

    @staticmethod
    def process(file_path: str, output_dir: str = None) -> ProcessedDocument:
        if not HAS_EXCEL:
            raise ImportError("openpyxl is not installed. Install it with: pip install openpyxl")

        start_time = time.time()
        file_path_obj = Path(file_path)

        try:
            wb = openpyxl.load_workbook(str(file_path_obj), read_only=True, data_only=True)
            content_parts = []

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                content_parts.append(f"\n[SHEET: {sheet_name}]")
                rows_text = []
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else '' for c in row]
                    if any(c.strip() for c in cells):
                        rows_text.append(" | ".join(cells))
                if rows_text:
                    content_parts.append("\n".join(rows_text))

            wb.close()
            content = "\n".join(content_parts)

            metadata = {
                'title': file_path_obj.stem,
                'author': 'Unknown',
                'word_count': len(content.split()),
                'char_count': len(content),
                'file_type': 'xlsx',
                'sheet_count': len(wb.sheetnames),
                'creation_date': datetime.fromtimestamp(file_path_obj.stat().st_ctime).isoformat()
            }

            return ProcessedDocument(
                content=content,
                metadata=metadata,
                doc_type='xlsx',
                file_path=str(file_path_obj),
                processing_time=time.time() - start_time
            )

        except Exception as e:
            raise Exception(f"Excel processing failed for {file_path}: {e}")


# ---------- MAIN MULTI-FORMAT PROCESSOR ----------



class MultiFormatDocumentProcessor:

    """Main processor that handles multiple file formats"""

    

    def __init__(self):

        self.processors = {

            'pdf': PDFProcessor,
            'pptx': PowerPointProcessor,
            'ppt': PowerPointProcessor,
            'docx': WordProcessor,
            'doc': WordProcessor,
            'text': TextProcessor,
            'txt': TextProcessor,
            'csv': TextProcessor,
            'md': TextProcessor,
            'xlsx': ExcelProcessor,
            'xls': ExcelProcessor,

        }

        

        self.processed_files: Dict[str, ProcessedDocument] = {}

    

    def process_file(self, file_path: str, output_dir: str = None) -> ProcessedDocument:

        """

        Process a file of any supported format

        

        Args:

            file_path: Path to the file to process

            output_dir: Directory where images should be saved (optional)

            

        Returns:

            ProcessedDocument with content and metadata

        """

        file_path = Path(file_path)

        

        if not file_path.exists():

            raise FileNotFoundError(f"File not found: {file_path}")

        

        # Detect file type

        category, specific_type = FileTypeDetector.get_file_type(str(file_path))

        

        if category == 'unknown':

            raise ValueError(f"Unsupported file type: {specific_type}")

        

        # Get appropriate processor

        processor_class = self.processors.get(category)

        if not processor_class:

            raise ValueError(f"No processor available for file type: {category}")

        

        print(f"📄 Processing {category.upper()} file: {file_path.name}")

        

        try:

            # Process the file

            processed_doc = processor_class.process(str(file_path), output_dir)

            

            # Add universal metadata

            processed_doc.metadata.update({

                'file_category': category,

                'file_size': file_path.stat().st_size,

                'file_modified': datetime.fromtimestamp(file_path.stat().st_mtime).isoformat(),

                'processor_version': '1.0.0'

            })

            

            # Cache the result

            file_hash = hashlib.md5(str(file_path).encode()).hexdigest()

            self.processed_files[file_hash] = processed_doc

            

            print(f"✅ Processed {file_path.name} in {processed_doc.processing_time:.2f}s")

            print(f"   Content length: {len(processed_doc.content)} chars")

            print(f"   Word count: {processed_doc.metadata.get('word_count', 0)}")

            

            return processed_doc

            

        except Exception as e:

            raise Exception(f"Failed to process {file_path.name}: {e}")

    

    def create_chunks(self, processed_doc: ProcessedDocument, 

                     chunk_size: int = 800, chunk_overlap: int = 150) -> List[Document]:

        """

        Split processed document into chunks for embedding

        

        Args:

            processed_doc: ProcessedDocument to chunk

            chunk_size: Target chunk size in characters

            chunk_overlap: Overlap between chunks

            

        Returns:

            List of Document chunks ready for embedding

        """

        if not processed_doc.content.strip():

            return []

        

        # Use appropriate splitter based on document type and content

        content_lower = processed_doc.content.lower()

        

        if processed_doc.doc_type in ['pptx']:

            # For presentations, prefer breaking on slides

            separators = ["SLIDE", "\n\n", "\n", ".", "!", "?", ",", " ", ""]

        else:

            # Default text splitting

            separators = ["\n\n", "\n", ".", "!", "?", ",", " ", ""]

        

        text_splitter = RecursiveCharacterTextSplitter(

            chunk_size=chunk_size,

            chunk_overlap=chunk_overlap,

            separators=separators

        )

        

        # Create initial document

        doc = Document(

            page_content=processed_doc.content,

            metadata=processed_doc.metadata.copy()

        )

        

        # Split into chunks

        chunks = text_splitter.split_documents([doc])

        

        # Enhance chunk metadata

        for i, chunk in enumerate(chunks):

            chunk.metadata.update({

                'chunk_index': i,

                'total_chunks': len(chunks),

                'chunk_size': len(chunk.page_content),

                'processing_method': 'unified_document_processor'

            })

            

            # Enhanced experiment detection for chunks

            chunk_text_lower = chunk.page_content.lower()

            experiment_indicators = ['experiment', 'lab', 'practical', 'exercise', 'tutorial', 'step', 'procedure', 'aim', 'objective']

            chunk.metadata['is_experiment'] = any(indicator in chunk_text_lower for indicator in experiment_indicators)

            

            # Extract experiment number if present in chunk

            import re

            experiment_match = re.search(r'experiment\s*(?:no\.?|number)?\s*:?\s*(\d+)', chunk_text_lower)

            if experiment_match:

                chunk.metadata['experiment_number'] = int(experiment_match.group(1))

            

            # Extract experiment aim if present in chunk

            aim_match = re.search(r'aim\s*:?\s*([^.\n]+)', chunk_text_lower)

            if aim_match:

                chunk.metadata['experiment_aim'] = aim_match.group(1).strip()

        

        # Cache chunks

        processed_doc.chunks = chunks

        

        print(f"✂️ Created {len(chunks)} chunks from {processed_doc.doc_type} document")

        

        return chunks

    

    def get_supported_formats(self) -> List[str]:

        """Get list of supported file formats"""

        return FileTypeDetector.get_supported_extensions()

    

    def is_supported(self, file_path: str) -> bool:

        """Check if file format is supported"""

        return FileTypeDetector.is_supported(file_path)



# ---------- UTILITY FUNCTIONS ----------



def get_document_name(file_path: str) -> str:

    """Extract document name without extension and replace spaces with underscores"""

    return Path(file_path).stem.replace(" ", "_")



def validate_document_path(file_path: str) -> bool:

    """Validate document file exists and is supported"""

    if not os.path.isfile(file_path):

        return False

    return FileTypeDetector.is_supported(file_path)



# ---------- EXPORT MAIN COMPONENTS ----------



__all__ = [

    'MultiFormatDocumentProcessor',

    'FileTypeDetector', 

    'ProcessedDocument',

    'PowerPointProcessor',

    'GraphProcessor',

    'EnhancedImageProcessor',

    'AdvancedTableProcessor',

    'get_document_name',

    'validate_document_path',

    'extract_text_from_image',

    'extract_names_from_content'

]

